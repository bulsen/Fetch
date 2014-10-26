#-*- encoding: utf-8 -*-
import sys
reload(sys)
sys.setdefaultencoding("utf-8")

#declarations
global parsed
parsed = ""
"""
    !!! feed only with pptx !!!

    for whom might need to convert ppts to pptx:
    libreoffice --headless --invisible --convert-to pptx filename.ppt

    at the end it gives you:

                pptxs in pptxdude/

                svgs in svgdude/name_of_that_file/ 


    note: 
        you can change the database and collection names, there is a abstraction at 48th line


"""


from pptx import Presentation

import sys, os 

merb= sys.argv[1]

if merb[:2] == "./":
    merb = merb[2:]

prs = Presentation(merb)

filename = merb[:-5]

"""	database json format:
   	
    { querry ,filename(scrambled), pagenumber, title(first 2 lines) }

"""

# database connection
from pymongo import MongoClient
client = MongoClient()

# you can change the database and collection names with
#           db = client.desired_database_name
#           collection = db.desired_collection_name

db = client.mikrodb
collection = db.motor


text_runs = []
n =0
for slide in prs.slides:
    n +=1
    for shape in slide.shapes:
        if not shape.has_textframe:
            continue
        for paragraph in shape.textframe.paragraphs:
            for run in paragraph.runs:
                # for title and stuff
                text_runs.append({ 'page': str(n) , 'text' : run.text})
                # database things
                posst = {"querry" : run.text , "filename" : filename, "pagenumber" : n}
                collection.insert(posst)
    #print "page count :" + str(n)

# this will extract pdfs for pdf2svg
os.system("libreoffice --headless --invisible --convert-to pdf %s.pptx"%filename)

""" SVG extractor """

os.system("mkdir svgdude/%s" %filename)
for temp in range(1,n+1):
	os.system("pdf2svg %s.pdf svgdude/%s/%s.svg %s" %(filename ,filename,temp ,temp))

os.system("mv %s.pptx pptxdude/%s.pptx"%(filename,filename))
os.system("rm -rf %s.pdf"%filename)
print "ok"
exit()
