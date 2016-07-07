# -*- encoding: utf-8 -*-
import sys
reload(sys)
sys.setdefaultencoding("utf-8")

import os, glob, random

crown = os.getcwd() +"/"

def rndnm(buff):
	array = "qwertyuopasdfghjklizxcvbnmQWERTYUIOASDFGHJKLZXCVBNM1234567890"
	ret = ""
	for i in range(buff):
		ret += array[random.randrange(0,len(array))]

	return ret 

slides = crown + "slides/"
pdfiles = crown + "pdfiles/"

for types in ["*.ppt","*.pptx"]:
	""" renames the specidfic types """
	for ppt in glob.glob(slides + types):
		dir = "/".join(ppt.split("/")[0:-1])+"/"
		os.rename(ppt,dir + rndnm(12)+types[1:])


for ptes in glob.glob(slides+"*.ppt"):
	""" pptx converter for ppts """
	os.system("libreoffice --headless --invisible  --convert-to pptx  --outdir "+ slides+" "+ptes)
	os.system("rm "+ptes)


for pxes in glob.glob(slides+"*.pptx"):
	os.system("libreoffice --headless --invisible  --convert-to pdf --outdir "+ pdfiles+" "+pxes)


from pptx import Presentation

import pymongo
cli = pymongo.MongoClient()
con = cli.fetchdatabase
mot = con.metadata


def extractor(filename):
	prs = Presentation(filename)

	text_runs = []
	n =0
	for slide in prs.slides:
	    n +=1
	    for shape in slide.shapes:
	        if not shape.has_text_frame:
	            continue
	        for paragraph in shape.text_frame.paragraphs:
	            for run in paragraph.runs:
	                # for title and stuff
	                text_runs.append({ 'page': str(n) , 'text' : run.text})
	                # database things
	                posst = {"querry" : run.text , "filename" : filename, "pagenumber" : n}
	                mot.insert(posst)
	    #print "page count :" + str(n)

	# this will extract pdfs for pdf2svg
	#os.system("libreoffice --headless --outdir %s --convert-to pdf %s"%(crown+"pdfiles/",filename))
	""" svgextracting """

	name = filename.split("/")[-1][:-5]
	os.system("mkdir svgdude/%s" %name)
	for temp in range(1,n+1):
		print "svg stuff"
		print filename, name
		cmd = "pdf2svg "+ pdfiles+name+'.pdf'+" "+crown+"svgdude/"+name+"/"+str(temp)+".svg "+str(temp)
		print cmd
		os.system(cmd)

	#os.system("rm -rf %s.pdf"%filename.split("/")[-1][:-5])

	print "ok - ",filename

for files in glob.glob(slides+"*.pptx"):
	extractor(files)
